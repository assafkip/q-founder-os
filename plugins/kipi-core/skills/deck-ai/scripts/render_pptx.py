#!/usr/bin/env python3
"""
Render source markdown + decisions.json to an editable PPTX.

Pipeline:
    source.md  ─┐
    decisions  ─┼─→ render_pptx.py ─→ deck.pptx  (every text frame editable)
    Unsplash   ─┘    (embeds images in file, no hotlinks)

Usage:
    python3 scripts/render_pptx.py <source.md> <decisions.json> <out.pptx>

Requires env var UNSPLASH_ACCESS_KEY (only when a slide has an image_keyword).
"""

import argparse
import io
import json
import os
import re
import sys
import urllib.request
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Local helper
sys.path.insert(0, str(Path(__file__).parent))
from fetch_images import resolve_url  # noqa: E402


# --- Parsing (same contract as v2 Slidev pipeline) ----------------------------

LAYOUTS = {
    "cover", "center", "statement", "two-cols", "image-right",
    "image-left", "quote", "end", "default",
}
IMAGE_LAYOUTS = {"cover", "image-right", "image-left"}

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Palette — neutral, serif-friendly. One clean look, no theme system in v3.
DARK = RGBColor(0x1A, 0x1A, 0x1A)
MUTED = RGBColor(0x55, 0x55, 0x55)
ACCENT = RGBColor(0xB8, 0x5C, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BG_TINT = RGBColor(0xFA, 0xF7, 0xF2)


def split_source(text):
    parts = re.split(r"\n-{3,}\s*\n", text.strip())
    return [p.strip() for p in parts if p.strip()]


def parse_slide(body):
    """Return (title, rest). Strip layout hint comments."""
    clean = re.sub(r"<!--\s*layout:\s*[a-z-]+\s*-->\s*", "", body).strip()
    title_match = re.match(r"#\s+(.+)", clean)
    title = title_match.group(1).strip() if title_match else ""
    rest = clean[title_match.end():].strip() if title_match else clean
    return title, rest


# --- Text helpers -------------------------------------------------------------

BOLD_RE = re.compile(r"\*\*(.+?)\*\*")
BULLET_RE = re.compile(r"^\s*[-*]\s+(.+)$")
NUMLIST_RE = re.compile(r"^\s*\d+\.\s+(.+)$")
BLOCKQUOTE_RE = re.compile(r"^\s*>\s+(.+)$")


def md_to_lines(rest):
    """Return a list of (kind, text) where kind ∈ {p, bullet, num, quote, h2}.
    Strips markdown marks Powerpoint can't render (keeps bold for run-level)."""
    lines = []
    for raw in rest.split("\n"):
        line = raw.rstrip()
        if not line.strip():
            if lines and lines[-1][0] != "blank":
                lines.append(("blank", ""))
            continue
        if line.startswith("## "):
            lines.append(("h2", line[3:].strip()))
            continue
        m = BULLET_RE.match(line)
        if m:
            lines.append(("bullet", m.group(1)))
            continue
        m = NUMLIST_RE.match(line)
        if m:
            lines.append(("num", m.group(1)))
            continue
        m = BLOCKQUOTE_RE.match(line)
        if m:
            lines.append(("quote", m.group(1)))
            continue
        lines.append(("p", line.strip()))
    while lines and lines[-1][0] == "blank":
        lines.pop()
    return lines


def write_runs(paragraph, text, font_size_pt, color=DARK):
    """Split **bold** runs, add them to `paragraph`."""
    parts = []
    pos = 0
    for m in BOLD_RE.finditer(text):
        if m.start() > pos:
            parts.append((text[pos:m.start()], False))
        parts.append((m.group(1), True))
        pos = m.end()
    if pos < len(text):
        parts.append((text[pos:], False))
    if not parts:
        parts = [(text, False)]

    first = True
    for chunk, is_bold in parts:
        if first:
            run = paragraph.runs[0] if paragraph.runs else paragraph.add_run()
            run.text = chunk
            first = False
        else:
            run = paragraph.add_run()
            run.text = chunk
        run.font.size = Pt(font_size_pt)
        run.font.color.rgb = color
        run.font.bold = is_bold


# --- Image helpers ------------------------------------------------------------

def fetch_image_bytes(url):
    with urllib.request.urlopen(url, timeout=30) as resp:
        return resp.read()


def image_stream(keyword, access_key, cache):
    """Resolve keyword → URL → bytes. Cache per-run."""
    if keyword in cache:
        return cache[keyword]
    try:
        url = resolve_url(keyword, access_key)
        if not url:
            cache[keyword] = None
            return None
        data = fetch_image_bytes(url)
    except Exception as err:
        print(f"  WARN: image failed for {keyword!r}: {err}", file=sys.stderr)
        cache[keyword] = None
        return None
    cache[keyword] = data
    return data


# --- Layout renderers ---------------------------------------------------------

def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])  # index 6 = blank


def add_textbox(slide, left, top, width, height, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    return tf


def add_background(slide, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    # Send to back
    spTree = shape._element.getparent()
    spTree.remove(shape._element)
    spTree.insert(2, shape._element)
    return shape


def render_body_lines(tf, lines, size=16, color=DARK):
    """Pour md lines into an existing text_frame. First paragraph reuses
    the frame's default; subsequent lines add new paragraphs."""
    first = True
    for kind, text in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        if kind == "blank":
            p.text = ""
            continue
        if kind == "h2":
            write_runs(p, text, size + 4, color=ACCENT)
            p.space_after = Pt(6)
            continue
        if kind == "bullet":
            p.level = 0
            write_runs(p, "•  " + text, size, color=color)
            p.space_after = Pt(4)
            continue
        if kind == "num":
            write_runs(p, text, size, color=color)
            p.space_after = Pt(4)
            continue
        if kind == "quote":
            write_runs(p, text, size + 2, color=MUTED)
            p.space_after = Pt(6)
            for r in p.runs:
                r.font.italic = True
            continue
        write_runs(p, text, size, color=color)
        p.space_after = Pt(6)


def layout_cover(prs, slide, title, rest, image_bytes):
    # Full-bleed image + dark overlay + big centered title.
    if image_bytes:
        slide.shapes.add_picture(io.BytesIO(image_bytes), 0, 0,
                                 width=SLIDE_W, height=SLIDE_H)
    overlay = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    overlay.line.fill.background()
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = RGBColor(0x0A, 0x0A, 0x0A)
    overlay.fill.transparency = 0  # python-pptx sets via _fill xml if needed
    # Best-effort transparency via raw XML (python-pptx doesn't expose it directly)
    _set_shape_transparency(overlay, 45)

    title_tf = add_textbox(
        slide, Inches(1.0), Inches(2.2), Inches(11.3), Inches(2.5),
        anchor=MSO_ANCHOR.MIDDLE,
    )
    title_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    write_runs(title_tf.paragraphs[0], title, 54, color=WHITE)
    title_tf.paragraphs[0].runs[0].font.bold = True

    subtitle = _extract_subtitle(rest)
    if subtitle:
        sub_tf = add_textbox(
            slide, Inches(1.0), Inches(4.9), Inches(11.3), Inches(1.0),
            anchor=MSO_ANCHOR.MIDDLE,
        )
        sub_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        write_runs(sub_tf.paragraphs[0], subtitle, 22, color=WHITE)


def layout_center(prs, slide, title, rest):
    add_background(slide, BG_TINT)
    lines = md_to_lines(rest)
    if title:
        title_tf = add_textbox(
            slide, Inches(1.0), Inches(1.8), Inches(11.3), Inches(1.2),
            anchor=MSO_ANCHOR.MIDDLE,
        )
        title_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        write_runs(title_tf.paragraphs[0], title, 40, color=DARK)
        title_tf.paragraphs[0].runs[0].font.bold = True
    body_tf = add_textbox(
        slide, Inches(1.5), Inches(3.2), Inches(10.3), Inches(3.5),
        anchor=MSO_ANCHOR.TOP,
    )
    for p in _each_pour(body_tf, lines, size=20):
        p.alignment = PP_ALIGN.CENTER


def layout_statement(prs, slide, title, rest):
    add_background(slide, DARK)
    text = (rest or title).strip()
    tf = add_textbox(
        slide, Inches(1.5), Inches(2.5), Inches(10.3), Inches(2.5),
        anchor=MSO_ANCHOR.MIDDLE,
    )
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    write_runs(tf.paragraphs[0], text, 44, color=WHITE)
    tf.paragraphs[0].runs[0].font.bold = True


def layout_quote(prs, slide, title, rest):
    add_background(slide, BG_TINT)
    lines = md_to_lines(rest)
    quote_text = " ".join(t for k, t in lines if k in ("quote", "p"))
    tf = add_textbox(
        slide, Inches(1.5), Inches(2.0), Inches(10.3), Inches(3.5),
        anchor=MSO_ANCHOR.MIDDLE,
    )
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    write_runs(p, f"\u201C{quote_text}\u201D", 34, color=DARK)
    for r in p.runs:
        r.font.italic = True
    if title:
        attr = add_textbox(
            slide, Inches(1.5), Inches(5.8), Inches(10.3), Inches(0.6),
            anchor=MSO_ANCHOR.TOP,
        )
        attr.paragraphs[0].alignment = PP_ALIGN.CENTER
        write_runs(attr.paragraphs[0], f"— {title}", 16, color=MUTED)


def layout_end(prs, slide, title, rest):
    add_background(slide, BG_TINT)
    title_tf = add_textbox(
        slide, Inches(1.0), Inches(2.5), Inches(11.3), Inches(1.2),
        anchor=MSO_ANCHOR.MIDDLE,
    )
    title_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    write_runs(title_tf.paragraphs[0], title or "Thanks", 44, color=ACCENT)
    title_tf.paragraphs[0].runs[0].font.bold = True
    if rest:
        body_tf = add_textbox(
            slide, Inches(1.5), Inches(4.2), Inches(10.3), Inches(2.5),
            anchor=MSO_ANCHOR.TOP,
        )
        for p in _each_pour(body_tf, md_to_lines(rest), size=18):
            p.alignment = PP_ALIGN.CENTER


def layout_default(prs, slide, title, rest):
    add_background(slide, WHITE)
    _title_bar(slide, title)
    body_tf = add_textbox(
        slide, Inches(0.8), Inches(1.7), Inches(11.7), Inches(5.5),
        anchor=MSO_ANCHOR.TOP,
    )
    _each_pour(body_tf, md_to_lines(rest), size=18)


def layout_two_cols(prs, slide, title, rest):
    add_background(slide, WHITE)
    _title_bar(slide, title)
    lines = md_to_lines(rest)
    mid = len(lines) // 2 or 1
    left_tf = add_textbox(
        slide, Inches(0.8), Inches(1.7), Inches(5.8), Inches(5.5),
        anchor=MSO_ANCHOR.TOP,
    )
    right_tf = add_textbox(
        slide, Inches(6.8), Inches(1.7), Inches(5.8), Inches(5.5),
        anchor=MSO_ANCHOR.TOP,
    )
    _each_pour(left_tf, lines[:mid], size=18)
    _each_pour(right_tf, lines[mid:], size=18)


def layout_image_side(prs, slide, title, rest, image_bytes, side):
    add_background(slide, WHITE)
    _title_bar(slide, title, width_in=7.5)
    body_left = Inches(0.8) if side == "left-text" else Inches(6.3)
    if image_bytes:
        img_left = Inches(7.0) if side == "left-text" else Inches(0.3)
        slide.shapes.add_picture(
            io.BytesIO(image_bytes), img_left, Inches(0.8),
            width=Inches(6.0), height=Inches(6.0),
        )
    body_tf = add_textbox(
        slide, body_left, Inches(1.7), Inches(6.2), Inches(5.5),
        anchor=MSO_ANCHOR.TOP,
    )
    _each_pour(body_tf, md_to_lines(rest), size=18)


# --- Small helpers ------------------------------------------------------------

def _title_bar(slide, title, width_in=12.0):
    if not title:
        return
    tf = add_textbox(slide, Inches(0.8), Inches(0.5), Inches(width_in),
                     Inches(1.0), anchor=MSO_ANCHOR.TOP)
    write_runs(tf.paragraphs[0], title, 28, color=DARK)
    tf.paragraphs[0].runs[0].font.bold = True


def _each_pour(tf, lines, size):
    """Render lines, yield the paragraphs so callers can tweak alignment."""
    first = True
    paragraphs = []
    for kind, text in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        paragraphs.append(p)
        if kind == "blank":
            p.text = ""
            continue
        if kind == "h2":
            write_runs(p, text, size + 4, color=ACCENT)
            p.space_after = Pt(6)
            continue
        if kind == "bullet":
            write_runs(p, "•  " + text, size, color=DARK)
            p.space_after = Pt(4)
            continue
        if kind == "num":
            write_runs(p, text, size, color=DARK)
            p.space_after = Pt(4)
            continue
        if kind == "quote":
            write_runs(p, text, size + 2, color=MUTED)
            for r in p.runs:
                r.font.italic = True
            continue
        write_runs(p, text, size, color=DARK)
        p.space_after = Pt(6)
    return paragraphs


def _extract_subtitle(rest):
    for line in rest.split("\n"):
        line = line.strip()
        if not line:
            continue
        if line.startswith("## "):
            return line[3:].strip()
        if line.startswith("#"):
            continue
        return line
    return ""


def _set_shape_transparency(shape, pct):
    """Set solid-fill alpha. pct = 0 opaque, 100 fully transparent."""
    from lxml import etree
    fill = shape.fill._xPr.find(
        "{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill"
    )
    if fill is None:
        return
    color = fill[0]
    alpha = etree.SubElement(
        color,
        "{http://schemas.openxmlformats.org/drawingml/2006/main}alpha",
    )
    alpha.set("val", str(int((100 - pct) * 1000)))


# --- Pipeline -----------------------------------------------------------------

def validate_decisions(data, slides):
    if not isinstance(data, list):
        raise ValueError("decisions must be a JSON array")
    if len(data) != len(slides):
        raise ValueError(
            f"decisions count {len(data)} != slides {len(slides)}"
        )
    out = []
    for i, d in enumerate(data):
        layout = d.get("layout")
        if layout not in LAYOUTS:
            raise ValueError(f"slide {i+1}: invalid layout {layout!r}")
        if i == 0 and layout != "cover":
            layout = "cover"
        keyword = d.get("image_keyword")
        if layout in IMAGE_LAYOUTS:
            if not isinstance(keyword, str) or not keyword.strip():
                raise ValueError(
                    f"slide {i+1} layout={layout} needs image_keyword"
                )
            keyword = " ".join(keyword.split()[:3])
        else:
            keyword = None
        out.append((layout, keyword))
    return out


def build_presentation(slides, picks, access_key):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    img_cache = {}

    for i, body in enumerate(slides):
        title, rest = parse_slide(body)
        layout, keyword = picks[i]
        slide = blank_slide(prs)

        img_bytes = None
        if keyword and access_key:
            img_bytes = image_stream(keyword, access_key, img_cache)

        if layout == "cover":
            layout_cover(prs, slide, title, rest, img_bytes)
        elif layout == "center":
            layout_center(prs, slide, title, rest)
        elif layout == "statement":
            layout_statement(prs, slide, title, rest)
        elif layout == "quote":
            layout_quote(prs, slide, title, rest)
        elif layout == "end":
            layout_end(prs, slide, title, rest)
        elif layout == "two-cols":
            layout_two_cols(prs, slide, title, rest)
        elif layout == "image-right":
            layout_image_side(prs, slide, title, rest, img_bytes,
                              side="left-text")
        elif layout == "image-left":
            layout_image_side(prs, slide, title, rest, img_bytes,
                              side="right-text")
        else:
            layout_default(prs, slide, title, rest)
    return prs


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("source", type=Path)
    parser.add_argument("decisions", type=Path)
    parser.add_argument("output", type=Path)
    args = parser.parse_args()

    if not args.source.exists():
        print(f"ERROR: source not found: {args.source}", file=sys.stderr)
        sys.exit(1)
    if not args.decisions.exists():
        print(f"ERROR: decisions not found: {args.decisions}", file=sys.stderr)
        sys.exit(1)

    slides = split_source(args.source.read_text())
    if not slides:
        print("ERROR: no slides in source", file=sys.stderr)
        sys.exit(1)

    try:
        picks = validate_decisions(
            json.loads(args.decisions.read_text()), slides
        )
    except ValueError as err:
        print(f"ERROR: {err}", file=sys.stderr)
        sys.exit(1)

    needs_key = any(kw for _, kw in picks)
    access_key = os.environ.get("UNSPLASH_ACCESS_KEY")
    if needs_key and not access_key:
        print("ERROR: UNSPLASH_ACCESS_KEY required (image slides present)",
              file=sys.stderr)
        sys.exit(1)

    prs = build_presentation(slides, picks, access_key)
    args.output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(args.output))
    print(f"WROTE: {args.output}  slides={len(slides)}")


if __name__ == "__main__":
    main()
